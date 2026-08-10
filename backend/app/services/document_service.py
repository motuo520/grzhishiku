import logging
import os
import shutil
import uuid
import re
from datetime import datetime
from typing import List, Optional

from sqlalchemy.orm import Session

from app.models.base import Document, User, KnowledgeUnit
from app.services import tag_service
from app.core.xss_sanitizer import sanitize_knowledge_input

logger = logging.getLogger(__name__)


def _get_safe_filename(filename: str) -> str:
    """Remove path components and sanitize filename."""
    base = os.path.basename(filename)
    base = re.sub(r"[^\w\-\.\u4e00-\u9fa5]", "_", base)
    return base


def _extract_text_from_txt(file_path: str) -> str:
    encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
    for enc in encodings:
        try:
            # strict 探测，避免 errors=replace 让 utf-8 假成功导致 GBK 中文乱码
            with open(file_path, "r", encoding=enc, errors="strict") as f:
                return f.read()
        except UnicodeDecodeError:
            continue
        except Exception:
            continue
    # 最终兜底：损坏文件用 replace 避免崩溃
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _extract_text_from_pdf(file_path: str) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        parts = []
        for page in reader.pages:
            try:
                text = page.extract_text()
                if text:
                    parts.append(text)
            except Exception:
                continue
        return "\n".join(parts)
    except Exception as e:
        raise RuntimeError(f"PDF 提取失败: {e}") from e


def _extract_text_from_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception as e:
        raise RuntimeError(f"DOCX 提取失败: {e}") from e


def _extract_text_from_xlsx(file_path: str) -> str:
    try:
        from contextlib import closing
        from openpyxl import load_workbook
        # closing 确保 read_only workbook 句柄被关闭，防批量解析泄漏
        with closing(load_workbook(file_path, data_only=True, read_only=True)) as wb:
            parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        parts.append(row_text)
            return "\n".join(parts)
    except Exception as e:
        raise RuntimeError(f"XLSX 提取失败: {e}") from e


def _extract_text_from_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as e:
        raise RuntimeError(f"PPTX 提取失败: {e}") from e


def _extract_text_from_html(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            html = f.read()
        text = re.sub(r"<[^>]+>", " ", html)
        text = re.sub(r"\s+", " ", text).strip()
        return text
    except Exception as e:
        raise RuntimeError(f"HTML 提取失败: {e}") from e


def extract_text(file_path: str, file_type: Optional[str] = None) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if file_type:
        file_type = file_type.lower()

    if ext in (".txt",) or file_type in ("text/plain",):
        return _extract_text_from_txt(file_path)
    if ext in (".md", ".markdown") or file_type in ("text/markdown",):
        return _extract_text_from_txt(file_path)
    if ext == ".pdf" or file_type in ("application/pdf",):
        return _extract_text_from_pdf(file_path)
    if ext == ".docx" or file_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ):
        return _extract_text_from_docx(file_path)
    if ext in (".xlsx", ".xls") or file_type in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        return _extract_text_from_xlsx(file_path)
    if ext == ".pptx" or file_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ):
        return _extract_text_from_pptx(file_path)
    if ext in (".html", ".htm") or file_type in ("text/html",):
        return _extract_text_from_html(file_path)

    raise ValueError(f"不支持的文件格式: {ext}")


ALLOWED_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".html", ".htm"
}


def create_document(
    db: Session,
    user: User,
    uploaded_file_path: str,
    original_name: str,
    file_type: Optional[str] = None,
    file_size: int = 0,
    title: Optional[str] = None,
) -> Document:
    ext = os.path.splitext(original_name)[1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"不支持的文件格式: {ext}")

    safe_name = _get_safe_filename(original_name)
    unique_name = f"{uuid.uuid4()}_{safe_name}"
    upload_dir = "uploads/documents"
    os.makedirs(upload_dir, exist_ok=True)
    dest_path = os.path.join(upload_dir, unique_name)
    shutil.move(uploaded_file_path, dest_path)

    doc = Document(
        id=str(uuid.uuid4()),
        user_id=user.id,
        title=title or safe_name,
        original_name=original_name,
        file_path=dest_path,
        file_size=file_size,
        file_type=file_type or ext,
        content_text=None,
        extraction_status="pending",
        doc_status="active",
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)

    # Extract text asynchronously-like (in same request for simplicity)
    try:
        content = extract_text(dest_path, file_type)
        doc.content_text = content
        doc.extraction_status = "success"
    except Exception as e:
        doc.extraction_status = "error"
        doc.extraction_error = str(e)

    doc.updated_at = datetime.now()
    db.commit()
    db.refresh(doc)
    return doc


def reextract_document(db: Session, user: User, document_id: str) -> Optional[Document]:
    doc = get_document(db, user, document_id)
    if not doc:
        return None
    try:
        content = extract_text(doc.file_path, doc.file_type)
        doc.content_text = content
        doc.extraction_status = "success"
        doc.extraction_error = None
    except Exception as e:
        doc.extraction_status = "error"
        doc.extraction_error = str(e)
    doc.updated_at = datetime.now()
    db.commit()
    db.refresh(doc)
    return doc


def list_documents(
    db: Session,
    user: User,
    file_type: Optional[str] = None,
    extraction_status: Optional[str] = None,
    q: Optional[str] = None,
    skip: int = 0,
    limit: int = 50,
) -> List[Document]:
    query = db.query(Document).filter(
        Document.user_id == user.id,
        Document.doc_status == "active"
    )
    if file_type:
        query = query.filter(Document.file_type.ilike(f"%{file_type}%"))
    if extraction_status:
        query = query.filter(Document.extraction_status == extraction_status)
    if q:
        search = f"%{q}%"
        query = query.filter(
            Document.title.ilike(search)
            | Document.original_name.ilike(search)
            | Document.content_text.ilike(search)
        )
    return query.order_by(Document.created_at.desc()).offset(skip).limit(limit).all()


def get_document(db: Session, user: User, document_id: str) -> Optional[Document]:
    return db.query(Document).filter(
        Document.id == document_id,
        Document.user_id == user.id,
        Document.doc_status == "active"
    ).first()


def delete_document(db: Session, user: User, document_id: str) -> bool:
    doc = get_document(db, user, document_id)
    if not doc:
        return False
    doc.doc_status = "deleted"
    db.commit()
    # Optionally delete file physically; keep for now to avoid accidental loss
    return True


def save_to_knowledge(db: Session, user: User, doc: Document, tag_ids: Optional[List[str]] = None) -> str:
    content = doc.content_text or ""
    safe_content, _, safe_title = sanitize_knowledge_input(
        content,
        None,
        doc.title or doc.original_name or "(无标题)"
    )

    unit = KnowledgeUnit(
        id=str(uuid.uuid4()),
        user_id=user.id,
        brain_side='network',
        content_raw=safe_content,
        content_type='document',
        source_url=f"/uploads/{doc.file_path}" if doc.file_path else None,
        source_title=safe_title,
        source_type='document',
        source_author=None,
        source_publish_date=None,
        verification_status='unverified',
        trust_level='tentative',
        verification_history='[]',
    )
    db.add(unit)
    db.commit()
    db.refresh(unit)

    if tag_ids:
        tag_service.set_tags_for(
            db,
            content_type=tag_service.CONTENT_TYPE_KNOWLEDGE,
            content_id=unit.id,
            user_id=user.id,
            tag_inputs=tag_ids,
        )
        db.commit()
        db.refresh(unit)

    try:
        from app.api.v1.endpoints.graph import auto_link_knowledge
        auto_link_knowledge(db, unit, user.id)
        db.commit()
    except Exception as e:
        logger.warning(f"Auto-link failed for document knowledge {unit.id}: {e}")

    doc.doc_status = "imported_to_knowledge"
    doc.knowledge_id = unit.id
    db.commit()

    return unit.id
